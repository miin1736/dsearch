"""
텍스트 추출 유틸리티

다양한 파일 형식(PDF, Word, Excel, PowerPoint, HTML 등)에서
텍스트 내용을 추출하는 유틸리티 모듈입니다.
각 파일 타입에 맞는 전용 추출기를 제공하며, 오류 처리와 로깅을 포함합니다.
"""

import os
import subprocess
from pathlib import Path
from typing import Optional, Dict, Any
import logging

logger = logging.getLogger("ds")


class TextExtractor:
    """\n    다양한 파일 형식에서 텍스트 내용을 추출하는 유틸리티 클래스.\n\n    PDF, Word, Excel, PowerPoint, HTML, RTF 등의 다양한 파일 형식에서\n    텍스트를 추출하고 전처리하여 검색 인덱싱에 사용할 수 있도록\n    텍스트 데이터를 제공합니다.\n    """

    def __init__(self):
        self.temp_dir = Path("/tmp/dsearch_extraction")
        self.temp_dir.mkdir(exist_ok=True)

    async def extract_text(self, file_path: str) -> Optional[str]:
        """\n        파일 타입에 따라 텍스트 내용을 추출합니다.\n\n        파일 확장자를 확인하여 적절한 추출기로 라우팅하고\n        각 파일 타입에 최적화된 방법으로 텍스트를 추출합니다.\n\n        Args:\n            file_path: 텍스트를 추출할 파일의 경로\n\n        Returns:\n            Optional[str]: 추출된 텍스트 내용 또는 None (추출 실패 시)\n\n        Raises:\n            Exception: 텍스트 추출 중 오류 발생 시\n        """
        try:
            path = Path(file_path)

            if not path.exists():
                logger.error(f"File not found: {file_path}")
                return None

            # Get file extension
            extension = path.suffix.lower()

            # Route to appropriate extractor
            if extension == ".pdf":
                return await self._extract_from_pdf(file_path)
            elif extension in [".doc", ".docx"]:
                return await self._extract_from_word(file_path)
            elif extension in [".xls", ".xlsx"]:
                return await self._extract_from_excel(file_path)
            elif extension in [".ppt", ".pptx"]:
                return await self._extract_from_powerpoint(file_path)
            elif extension in [".txt", ".md", ".csv"]:
                return await self._extract_from_text(file_path)
            elif extension in [".html", ".htm"]:
                return await self._extract_from_html(file_path)
            elif extension in [".rtf"]:
                return await self._extract_from_rtf(file_path)
            else:
                logger.warning(f"Unsupported file type: {extension}")
                return None

        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {e}")
            return None

    async def _extract_from_pdf(self, file_path: str) -> Optional[str]:
        """Extract text from PDF files."""
        try:
            # Try PyMuPDF first (more reliable)
            try:
                import fitz  # PyMuPDF

                doc = fitz.open(file_path)
                text_content = ""

                for page_num in range(doc.page_count):
                    page = doc[page_num]
                    text_content += page.get_text()

                doc.close()
                return text_content.strip()

            except ImportError:
                # Fall back to pdfplumber or PyPDF2
                try:
                    import pdfplumber

                    with pdfplumber.open(file_path) as pdf:
                        text_content = ""
                        for page in pdf.pages:
                            page_text = page.extract_text()
                            if page_text:
                                text_content += page_text + "\n"

                    return text_content.strip()

                except ImportError:
                    # Fall back to PyPDF2
                    try:
                        import PyPDF2

                        with open(file_path, 'rb') as file:
                            reader = PyPDF2.PdfReader(file)
                            text_content = ""

                            for page in reader.pages:
                                text_content += page.extract_text() + "\n"

                        return text_content.strip()

                    except ImportError:
                        # Fall back to pdftotext command
                        return await self._extract_pdf_with_pdftotext(file_path)

        except Exception as e:
            logger.error(f"Error extracting from PDF {file_path}: {e}")
            return None

    async def _extract_pdf_with_pdftotext(self, file_path: str) -> Optional[str]:
        """Extract PDF text using pdftotext command-line tool."""
        try:
            result = subprocess.run(
                ["pdftotext", "-layout", file_path, "-"],
                capture_output=True,
                text=True,
                timeout=30
            )

            if result.returncode == 0:
                return result.stdout.strip()
            else:
                logger.error(f"pdftotext error: {result.stderr}")
                return None

        except subprocess.TimeoutExpired:
            logger.error(f"PDF text extraction timeout: {file_path}")
            return None
        except FileNotFoundError:
            logger.error("pdftotext command not found. Install poppler-utils.")
            return None
        except Exception as e:
            logger.error(f"Error with pdftotext: {e}")
            return None

    async def _extract_from_word(self, file_path: str) -> Optional[str]:
        """Extract text from Word documents."""
        try:
            # Try python-docx for .docx files
            if file_path.endswith(".docx"):
                try:
                    from docx import Document

                    doc = Document(file_path)
                    text_content = ""

                    for paragraph in doc.paragraphs:
                        text_content += paragraph.text + "\n"

                    # Extract from tables
                    for table in doc.tables:
                        for row in table.rows:
                            for cell in row.cells:
                                text_content += cell.text + " "
                        text_content += "\n"

                    return text_content.strip()

                except ImportError:
                    pass

            # Fall back to LibreOffice
            return await self._extract_with_libreoffice(file_path)

        except Exception as e:
            logger.error(f"Error extracting from Word document {file_path}: {e}")
            return None

    async def _extract_from_excel(self, file_path: str) -> Optional[str]:
        """Extract text from Excel files."""
        try:
            # Try openpyxl for .xlsx files
            if file_path.endswith(".xlsx"):
                try:
                    from openpyxl import load_workbook

                    workbook = load_workbook(file_path, data_only=True)
                    text_content = ""

                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text_content += f"Sheet: {sheet_name}\n"

                        for row in sheet.iter_rows(values_only=True):
                            row_text = "\t".join([str(cell) if cell is not None else "" for cell in row])
                            if row_text.strip():
                                text_content += row_text + "\n"

                    return text_content.strip()

                except ImportError:
                    pass

            # Fall back to LibreOffice
            return await self._extract_with_libreoffice(file_path)

        except Exception as e:
            logger.error(f"Error extracting from Excel file {file_path}: {e}")
            return None

    async def _extract_from_powerpoint(self, file_path: str) -> Optional[str]:
        """Extract text from PowerPoint files."""
        try:
            # Try python-pptx for .pptx files
            if file_path.endswith(".pptx"):
                try:
                    from pptx import Presentation

                    prs = Presentation(file_path)
                    text_content = ""

                    for slide_num, slide in enumerate(prs.slides, 1):
                        text_content += f"Slide {slide_num}:\n"

                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_content += shape.text + "\n"

                        text_content += "\n"

                    return text_content.strip()

                except ImportError:
                    pass

            # Fall back to LibreOffice
            return await self._extract_with_libreoffice(file_path)

        except Exception as e:
            logger.error(f"Error extracting from PowerPoint file {file_path}: {e}")
            return None

    async def _extract_from_text(self, file_path: str) -> Optional[str]:
        """Extract text from plain text files."""
        try:
            # Try different encodings
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin1']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        return file.read()
                except UnicodeDecodeError:
                    continue

            logger.warning(f"Could not decode text file {file_path} with any encoding")
            return None

        except Exception as e:
            logger.error(f"Error extracting from text file {file_path}: {e}")
            return None

    async def _extract_from_html(self, file_path: str) -> Optional[str]:
        """Extract text from HTML files."""
        try:
            from bs4 import BeautifulSoup

            # Try different encodings
            encodings = ['utf-8', 'cp949', 'euc-kr', 'latin1']

            for encoding in encodings:
                try:
                    with open(file_path, 'r', encoding=encoding) as file:
                        html_content = file.read()

                    soup = BeautifulSoup(html_content, 'html.parser')

                    # Remove script and style elements
                    for script in soup(["script", "style"]):
                        script.decompose()

                    # Get text content
                    text = soup.get_text()

                    # Clean up text
                    lines = (line.strip() for line in text.splitlines())
                    chunks = (phrase.strip() for line in lines for phrase in line.split("  "))
                    text = ' '.join(chunk for chunk in chunks if chunk)

                    return text

                except UnicodeDecodeError:
                    continue

            logger.warning(f"Could not decode HTML file {file_path} with any encoding")
            return None

        except ImportError:
            logger.error("BeautifulSoup not available for HTML extraction")
            return None
        except Exception as e:
            logger.error(f"Error extracting from HTML file {file_path}: {e}")
            return None

    async def _extract_from_rtf(self, file_path: str) -> Optional[str]:
        """Extract text from RTF files."""
        try:
            # Use LibreOffice for RTF extraction
            return await self._extract_with_libreoffice(file_path)

        except Exception as e:
            logger.error(f"Error extracting from RTF file {file_path}: {e}")
            return None

    async def _extract_with_libreoffice(self, file_path: str) -> Optional[str]:
        """Extract text using LibreOffice headless mode."""
        try:
            # Create temporary output file
            temp_output = self.temp_dir / f"extracted_{Path(file_path).stem}.txt"

            # Run LibreOffice conversion
            result = subprocess.run([
                "libreoffice",
                "--headless",
                "--convert-to", "txt:Text",
                "--outdir", str(self.temp_dir),
                file_path
            ], capture_output=True, text=True, timeout=60)

            if result.returncode == 0:
                # Read the converted text file
                expected_output = self.temp_dir / f"{Path(file_path).stem}.txt"

                if expected_output.exists():
                    with open(expected_output, 'r', encoding='utf-8') as file:
                        text_content = file.read()

                    # Clean up temporary file
                    try:
                        expected_output.unlink()
                    except:
                        pass

                    return text_content.strip()

            logger.error(f"LibreOffice conversion failed: {result.stderr}")
            return None

        except subprocess.TimeoutExpired:
            logger.error(f"LibreOffice extraction timeout: {file_path}")
            return None
        except FileNotFoundError:
            logger.error("LibreOffice command not found")
            return None
        except Exception as e:
            logger.error(f"Error with LibreOffice extraction: {e}")
            return None

    async def extract_with_ocr(self, file_path: str, language: str = "kor+eng") -> Optional[str]:
        """Extract text using OCR (for scanned documents or images)."""
        try:
            # Use Tesseract OCR
            result = subprocess.run([
                "tesseract",
                file_path,
                "stdout",
                "-l", language
            ], capture_output=True, text=True, timeout=120)

            if result.returncode == 0:
                return result.stdout.strip()
            else:
                logger.error(f"Tesseract OCR error: {result.stderr}")
                return None

        except subprocess.TimeoutExpired:
            logger.error(f"OCR extraction timeout: {file_path}")
            return None
        except FileNotFoundError:
            logger.error("Tesseract OCR not found. Install tesseract-ocr.")
            return None
        except Exception as e:
            logger.error(f"Error with OCR extraction: {e}")
            return None

    def get_extraction_info(self, file_path: str) -> Dict[str, Any]:
        """Get information about text extraction capabilities for a file."""
        try:
            path = Path(file_path)
            extension = path.suffix.lower()

            extraction_methods = {
                ".pdf": ["PyMuPDF", "pdfplumber", "PyPDF2", "pdftotext", "OCR"],
                ".docx": ["python-docx", "LibreOffice"],
                ".doc": ["LibreOffice"],
                ".xlsx": ["openpyxl", "LibreOffice"],
                ".xls": ["LibreOffice"],
                ".pptx": ["python-pptx", "LibreOffice"],
                ".ppt": ["LibreOffice"],
                ".txt": ["Direct read"],
                ".html": ["BeautifulSoup"],
                ".htm": ["BeautifulSoup"],
                ".rtf": ["LibreOffice"]
            }

            return {
                "supported": extension in extraction_methods,
                "methods": extraction_methods.get(extension, []),
                "requires_ocr": extension in [".png", ".jpg", ".jpeg", ".tiff", ".bmp"],
                "file_type": extension
            }

        except Exception as e:
            logger.error(f"Error getting extraction info for {file_path}: {e}")
            return {"supported": False, "error": str(e)}

    async def cleanup_temp_files(self) -> int:
        """Clean up temporary extraction files."""
        try:
            deleted_count = 0

            if self.temp_dir.exists():
                for file_path in self.temp_dir.iterdir():
                    if file_path.is_file():
                        try:
                            file_path.unlink()
                            deleted_count += 1
                        except Exception as e:
                            logger.error(f"Error deleting temp file {file_path}: {e}")

            logger.info(f"Cleaned up {deleted_count} extraction temp files")
            return deleted_count

        except Exception as e:
            logger.error(f"Error during extraction temp file cleanup: {e}")
            return 0